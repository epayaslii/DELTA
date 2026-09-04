"""Generate the progress deck. python-pptx, 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x1F, 0x6F, 0x4D)      # herb green
FLAG = RGBColor(0xB4, 0x53, 0x09)        # amber
INK = RGBColor(0x22, 0x26, 0x2B)
MUTED = RGBColor(0x5B, 0x64, 0x70)
LIGHT = RGBColor(0xEC, 0xEE, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

FONT = "Calibri"


def _tf(box):
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_bar(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.95))
    tf = _tf(box)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.name = FONT; r2.font.color.rgb = MUTED
    # rule
    ln = slide.shapes.add_shape(1, Inches(0.55), Inches(1.45), Inches(12.2), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    return slide


def bullets(title, items, sub=None, body_top=1.75, size=17):
    """items: list of (level, text) or plain str (level 0). '' text -> spacer."""
    s = prs.slides.add_slide(BLANK); bg(s)
    title_bar(s, title, sub)
    box = s.shapes.add_textbox(Inches(0.7), Inches(body_top), Inches(12.0), Inches(5.4))
    tf = _tf(box)
    first = True
    for it in items:
        lvl, txt = it if isinstance(it, tuple) else (0, it)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(6)
        if txt == "":
            p.add_run().text = " "; p.space_after = Pt(4); continue
        bullet = "▸ " if lvl == 0 else ("– " if lvl == 1 else "· ")
        r = p.add_run(); r.text = bullet + txt
        r.font.name = FONT
        r.font.size = Pt(size if lvl == 0 else size - 2 if lvl == 1 else size - 3)
        r.font.color.rgb = INK if lvl == 0 else MUTED
        if lvl == 0:
            r.font.bold = False
    return s


def table_slide(title, headers, rows, sub=None, col_widths=None, font=12):
    s = prs.slides.add_slide(BLANK); bg(s)
    title_bar(s, title, sub)
    nr, nc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nr, nc, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.4 * nr)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(font + 1)
        pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.name = FONT
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            for pr in c.text_frame.paragraphs:
                for rn in pr.runs:
                    rn.font.size = Pt(font); rn.font.name = FONT; rn.font.color.rgb = INK
    return s


def section(title, kicker=None):
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.6))
    tf = _tf(box)
    if kicker:
        p = tf.paragraphs[0]; r = p.add_run(); r.text = kicker.upper()
        r.font.size = Pt(14); r.font.name = FONT; r.font.color.rgb = RGBColor(0x9D, 0xD3, 0xBE)
        r.font.bold = True
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(34); r2.font.bold = True; r2.font.name = FONT; r2.font.color.rgb = WHITE
    return s


# ------------------------------------------------------------------ TITLE
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
band = s.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(2.5))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
box = s.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.5))
tf = _tf(box)
r = tf.paragraphs[0].add_run()
r.text = "Temporal Alignment for DELTA"
r.font.size = Pt(40); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
p = tf.add_paragraph(); r = p.add_run()
r.text = "Progress, findings, and directions for a VLM-based approach"
r.font.size = Pt(18); r.font.name = FONT; r.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE3)
box = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.5))
tf = _tf(box)
for i, line in enumerate([
    "Eliz Payasli — research internship, IRI / UPC",
    "Base work: DELTA — Dense Long-Term Action Anticipation from Procedural Transcripts",
    "Repository: github.com/epayaslii/DELTA",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    rn = p.add_run(); rn.text = line
    rn.font.size = Pt(15 if i < 2 else 13); rn.font.name = FONT
    rn.font.color.rgb = INK if i < 2 else MUTED
    if i == 0:
        rn.font.bold = True

# ------------------------------------------------------------------ 1. WHERE WE ARE
bullets("Where we are — one slide", [
    (0, "Goal: transcript-only Dense Long-Term Action Anticipation (DLTA); focus = the "
        "Temporal Alignment (TA) component, on 50Salads."),
    (0, "Direction confirmed (your guidance): the TA must use a VLM, NOT segmentation the way ATBA/HAL do."),
    (0, "Done: research repo + infrastructure, 50Salads analysis, a VLM-direct aligner prototype "
        "(21 tests); deep-read of the DELTA code (received); analysed the CVPR'25/26 VLM-alignment "
        "literature (HAL, CVA, MASRA, TAN, StepFormer, OVTAS, MLLM4WTAL)."),
    (0, "Workstream split:  MASRA track  ·  CVA track.  Both are VLM video-text alignment."),
    (1, "HAL was analysed then dropped — it is ATBA + a regulariser, still segmentation-based."),
    (0, "Blocker: raw 50Salads video for VLM feature extraction (official host down)."),
    (0, "Next: start building MASRA + its VLM tomorrow."),
], sub="Internship progress review  ·  pivot: HAL → MASRA (VLM required)")

# ------------------------------------------------------------------ 2. PROBLEM RECAP
bullets("The problem, briefly", [
    (0, "DLTA: from an observed video prefix, predict the future actions, their order, and durations — densely."),
    (0, "DELTA learns this from transcripts only (ordered action list, no timestamps, no frame labels)."),
    (0, "Its Temporal Alignment module turns the transcript into dense pseudo-labels Y* that supervise "
        "segmentation and the anticipation decoder."),
    (0, "DELTA's TA follows ATBA (CVPR'24). On 50Salads DELTA trails fully-supervised methods:"),
    (1, "avg MoC 20.9 (DELTA)  vs  ~25.9 (FUTR)  vs  ~28.4 (ActFusion)."),
    (0, "If the pseudo-labels are poor, everything downstream is poor. So TA quality is the lever."),
], sub="DLTA · transcript-only · why TA matters")

# ------------------------------------------------------------------ 3. FRAMING
bullets("The research framing", [
    (0, "ATBA / HAL — \"alignment THROUGH segmentation\"  (ruled out by your no-segmentation guidance):"),
    (1, "frozen I3D  →  trained frame classifier  →  posteriors P  →  boundary/transition scores  →  DP  →  Y*"),
    (1, "the alignment can be no better than that classifier — and on 50Salads (fixed camera, "
        "near-duplicate cut_* / add_* actions) it is a near-worst case."),
    (0, "Our direction — \"alignment THROUGH semantic matching\":"),
    (1, "frozen VLM:  s(n, t) = sim( text(action_n), video_t )  →  order-preserving alignment  →  Y*"),
    (1, "no frame classifier; the fine-grained vocabulary is disambiguated by the noun, from step 0."),
    (0, "Two CVPR'25/26 papers do exactly this for video temporal grounding, and we adapt them:"),
    (1, "CVA (CVPR'26) — the aligner: hierarchical encoder + boundary-contrastive loss."),
    (1, "MASRA (2026) — the language regulariser: align a text relation-matrix to the video similarity "
        "matrix; MLLM used at training only, discarded at inference.  (Eliz)"),
], sub="the TA must use a VLM  ·  CVA + MASRA as the references")

# ------------------------------------------------------------------ 4. THE REPO
table_slide("What's built — the repository",
    ["Module", "What it does", "Status"],
    [
        ["delta.data", "50Salads / Breakfast conventions, transcripts, splits, dataset stats", "done + tests"],
        ["delta.features", "Frozen VLM feature extraction — VideoLLaMA3 / SigLIP2 / DINOv2 backbones, text encoder, extraction CLI", "done (untested on video)"],
        ["delta.align", "similarity matrix · order-preserving DP + soft alignment · segmentation/alignment metrics", "done + tests"],
        ["delta.viz", "segmentation-timeline plotting (GT vs pseudo-labels vs prediction)", "done"],
        ["docs/", "TA reference · 50Salads analysis · DELTA-code deep-read (loss assembly, VLM plug-in points) · HAL & CVA baselines · approach + lit review", "8 documents"],
        ["tests/", "CPU-only unit tests — synthetic + real-transcript validation", "21 passing"],
    ],
    sub="github.com/epayaslii/DELTA  ·  Python 3.11 / uv",
    col_widths=[2.1, 8.0, 2.0], font=11)

# ------------------------------------------------------------------ 5. PROGRESS: INFRA + DATA
bullets("Progress 1 — infrastructure & data", [
    (0, "Feature-extraction pipeline: pluggable frozen backbones (VL3-SigLIP = VideoLLaMA3 vision tower, "
        "SigLIP2, DINOv2), outputs (D,T) features aligned to the ground-truth frame grid — drop-in for I3D."),
    (0, "50Salads benchmark bundle downloaded & verified (HuggingFace dinggd/50salads):"),
    (1, "50 videos, 19 classes (17 actions + start/end), I3D 2048-d, confirmed 30 fps, ~20 segments/video."),
    (0, "Raw 50Salads videos: the official host (Dundee) is DOWN (NXDOMAIN everywhere); no public mirror."),
    (1, "Blocks VLM feature extraction. Need a copy from the lab — see questions at the end."),
    (0, "Environment: local Mac has no GPU/torch → model runs go to the UPC cluster."),
], sub="pipeline · dataset · the video blocker")

# ------------------------------------------------------------------ 6. PROGRESS: 50S ANALYSIS
table_slide("Progress 2 — why 50Salads is hard (measured)",
    ["Finding", "Number", "Implication"],
    [
        ["Naive uniform alignment (zero visual evidence)", "MoC 0.34", "the floor any TA must beat"],
        ["I3D consecutive-frame distance at GT boundaries\nvs at random frames", "1.11×", "class-agnostic boundary signal is\nnearly absent → ATBA's candidate\nstep misses true transitions"],
        ["Next-action entropy (transcript order)", "~1.9 bits\n(~3–4 successors)", "order alone doesn't pin timing —\nvisual evidence must carry it"],
        ["Per-class segment-duration variation", "CV 0.5–0.8", "duration head has weak signal\n(paper: +0.2 MoC on 50S vs +3.3 BF)"],
    ],
    sub="Stage 1 dataset analysis  ·  notebook + docs/50salads-notes.md",
    col_widths=[4.6, 2.3, 5.2], font=11)

# ------------------------------------------------------------------ 7. PROGRESS: VLM ALIGNER
bullets("Progress 3 — VLM-direct aligner prototype", [
    (0, "delta.align.similarity — transcript × frame cosine similarity from frozen VLM embeddings "
        "(model-agnostic; takes pre-extracted embeddings)."),
    (0, "delta.align.ta — two solvers:"),
    (1, "align_dp: hard order-preserving DP (monotone, every transcript action covered, transition penalty)."),
    (1, "align_soft: entropic forward–backward → per-frame posteriors AND P(boundary_r = t) distributions "
        "(for uncertainty-aware targets later)."),
    (0, "delta.align.evaluate — swap-in similarity providers (naive / oracle / VLM), score Y* vs held-out GT."),
    (0, "Validated align_dp on real 50Salads transcripts (up to 26 segments, repeated classes): "
        "near-perfect recovery from clean/noisy blocky similarity."),
    (0, "Missing piece: the real VLM similarity matrix → needs frame features → needs raw video."),
], sub="the method, minus the VLM  ·  21 tests")

# ------------------------------------------------------------------ 8. PROGRESS: DELTA CODE
bullets("Progress 4 — the DELTA code (\"WLTA\")", [
    (0, "Received from the group; analysed in docs/delta-code.md. WLTA = Weakly-supervised Long-Term Anticipation."),
    (0, "Built on the CLOT / ASOT codebase — so it already contains BOTH alignment mechanisms:"),
    (1, "--model_type atba  → the ATBA boundary detector (src/atba_loss.py)"),
    (1, "--model_type wclot → ASOT optimal transport (src/asot.py, gsw.py)"),
    (1, "plus CTC, TSM smoothing, cross-modal attention, linear-chain CRF, DistilBERT, full LTA eval."),
    (0, "The \"situation\": README is CLOT's; run scripts call a missing train_edit_elena9sept.py; "
        "hardcoded paths; needs Linux + CUDA + conda + wandb."),
    (0, "Implication: our contribution narrows cleanly — swap the frame-classifier posteriors / OT cost "
        "matrix for a frozen-VLM similarity; reuse everything downstream. And we can run the real baseline."),
], sub="what it is · how it runs · what it means for us")

# ------------------------------------------------------------------ 9. LITERATURE
table_slide("Progress 5 — VLM-alignment literature (analysed)",
    ["Paper", "Venue", "What it is", "Role for us"],
    [
        ["ATBA", "CVPR'24", "the classifier→DP alignment DELTA's TA follows", "the baseline we replace"],
        ["HAL", "CVPR'26", "= ATBA + a VAE regulariser; +2–3 MoF; segmentation-based; skips 50Salads", "analysed, then DROPPED (no-segmentation)"],
        ["CVA", "CVPR'26", "VLM video-text alignment for grounding; hierarchical encoder + boundary-contrastive loss; SOTA", "the VLM aligner"],
        ["MASRA", "2026", "MLLM-assisted alignment: align a text relation-matrix to the video similarity matrix; MLLM train-only", "Eliz — the language regulariser"],
        ["TAN / StepFormer", "CVPR'22/'23", "the genuine transcript/sequence→video aligners (weak/self-sup)", "the alignment mechanism references"],
    ],
    sub="no single paper = {weak supervision} × {VLM alignment} × {long-term anticipation}  →  that is the contribution",
    col_widths=[1.9, 1.0, 6.3, 2.9], font=9.5)

# ------------------------------------------------------------------ 10. MASRA + CVA
bullets("Progress 6 — the two VLM-alignment references", [
    (0, "CVA (CVPR'26) — Context-aware Video-text Alignment.  Task: video temporal grounding "
        "(query → span), CLIP+SlowFast, SOTA.  Contribution = the aligner:"),
    (1, "CTE hierarchical encoder (windowed self-attn + learnable queries + bidirectional cross-attn)"),
    (1, "CBD boundary-contrastive loss + QCD query-aware augmentation."),
    (0, "MASRA (2026) — MLLM-Assisted Semantic-Relational Consistent Alignment.  Same task, "
        "CLIP + an MLLM.  Contribution = a training-time language regulariser:"),
    (1, "LRCA — align a text relation-matrix (from MLLM captions) with the video's similarity matrix"),
    (1, "ESTA — align pooled temporal context with action/event semantics"),
    (1, "MLLM used ONLY at training, discarded at inference  (= DELTA's philosophy)."),
    (0, "Both: grounding, supervised, not procedural — we adapt them to transcript-supervised 50Salads."),
], sub="CVA = the aligner   ·   MASRA = the language regulariser")

# ------------------------------------------------------------------ 11. THE SPLIT
bullets("Scoping — the workstream split", [
    (0, "Supervisor's direction: the TA must use a VLM; no segmentation-based approach."),
    (0, ""),
    (0, "CVA track:  build the VLM aligner — CTE encoder + CBD boundary-contrastive loss."),
    (0, "Eliz → MASRA:  the training-time language regulariser — LRCA / ESTA — that shapes the "
        "video↔transcript similarity, then an order-preserving alignment reads Y* off it."),
    (0, ""),
    (0, "Structurally: CVA produces the alignment; MASRA (like HAL did on the segmentation side) "
        "is an auxiliary training signal that improves it — measured on TA metrics, then downstream MoC."),
    (0, "VLM for both: InternVideo2 (video-native) instead of CVA/MASRA's SlowFast + frame-CLIP."),
], sub="MASRA track   ·   CVA track   ·   both VLM")

# ------------------------------------------------------------------ 12. TAKEAWAYS
bullets("Key takeaways", [
    (0, "\"Alignment through segmentation\" (ATBA / HAL) is a trained classifier that 50Salads breaks — "
        "and the supervisor has ruled it out."),
    (0, "The move: replace the classifier posteriors with a frozen-VLM transcript×frame similarity, "
        "then an order-preserving alignment."),
    (0, "DELTA already has optimal-transport alignment (ASOT, --model_type wclot) — the VLM swap is "
        "≈ one line in the cost matrix; CTC / CRF / LTA decoder / eval are all reused."),
    (0, "CVA + MASRA (CVPR'26 / 2026) do VLM video-text alignment for grounding; we adapt them to "
        "transcript-supervised dense anticipation — which nobody has done."),
    (0, "HAL analysed and dropped; kept only as a baseline number."),
], sub="what the analysis established")

# ------------------------------------------------------------------ 13. BLOCKERS
table_slide("Blockers & risks",
    ["Item", "Impact", "Mitigation"],
    [
        ["Raw 50Salads video unavailable", "blocks VLM feature extraction (the core of both tracks)", "chase a lab copy; meanwhile reproduce MASRA/CVA on TACoS (cooking VTG benchmark)"],
        ["No local GPU / torch", "no model runs on the Mac", "UPC cluster; one conda env"],
        ["MoF ≠ MoC", "a method can win on segmentation MoF and not help DELTA's anticipation MoC", "always measure pseudo-label MoC + boundary offset, then downstream MoC"],
        ["DELTA code is research-grade", "reproduction friction (missing script, paths, wandb)", "documented in docs/delta-code.md; map run scripts onto train.py"],
        ["MASRA/CVA are grounding + supervised", "not a drop-in — architecture + losses must be adapted", "keep only the transferable pieces (LRCA/ESTA; CTE/CBD); wire into DELTA's ASOT + decoder"],
    ],
    sub="known before we commit compute",
    col_widths=[3.0, 4.3, 4.8], font=10.5)

# ------------------------------------------------------------------ 14. SECTION
section("How we proceed with the TA", kicker="For discussion")

# ------------------------------------------------------------------ 15. SCOPE
bullets("Scope — 50Salads only", [
    (0, "No Breakfast, no segmentation-based approach, no HAL workstream (kept as a baseline number)."),
    (0, "One dataset, evaluated two ways — TA metrics gate the DLTA metrics:"),
    (1, "TA metrics:  MoF / MoC / Edit / F1@{10,25,50} / median per-transition boundary offset — on Y* vs held-out GT."),
    (1, "DLTA metrics:  Obs{20,30}% × Pred{10,20,30,50}% MoC — the DELTA anticipation grid."),
    (0, "Method = frozen VLM (InternVideo2) transcript×frame similarity  →  MASRA-style regularisation "
        "(LRCA/ESTA)  →  order-preserving alignment (ASOT, already in DELTA)  →  Y*  →  DELTA decoder."),
    (0, "Baselines:  naive-uniform 0.34 MoC  ·  ATBA-in-DELTA (--model_type atba)  ·  ASOT-in-DELTA "
        "(--model_type wclot)  ·  HAL number (from its paper)."),
], sub="VLM-only · TA metrics → DLTA metrics · one dataset")

# ------------------------------------------------------------------ 16. ELIZ / MASRA PLAN
table_slide("My workstream — MASRA (M1–M7)",
    ["#", "Step", "Output"],
    [
        ["M1", "Get the MASRA code + a VTG base model; reproduce on TACoS (cooking benchmark — no 50S video needed yet)", "understand LRCA / ESTA / DAI"],
        ["M2", "Add an internvideo2 backbone to delta.features.backbones; build s(transcript × frame) for 50S", "s (N×T)   [needs raw video]"],
        ["M3", "Adapt LRCA — align a transcript relation-matrix (order structure) with the video similarity matrix", "language regulariser loss"],
        ["M4", "Adapt ESTA — align pooled temporal context with the 17 action-name semantics", "semantic alignment loss"],
        ["M5", "Order-preserving alignment (ASOT / delta.align.ta) reads Y* off the regularised similarity", "Y*"],
        ["M6", "Score Y* vs ATBA-in-DELTA / ASOT-in-DELTA / naive floor — TA metrics", "does VLM+MASRA beat the baselines?"],
        ["M7", "Feed best Y* into DELTA's decoder (CTC/CRF/duration unchanged) → Obs%/Pred% MoC", "the downstream result"],
    ],
    sub="frozen VLM similarity + MASRA-style language regularisation → alignment → DELTA.  M1 starts tomorrow.",
    col_widths=[0.5, 8.6, 3.0], font=9)

# ------------------------------------------------------------------ 17. THE FULL METHOD
bullets("The full method — where CVA + MASRA land", [
    (0, "s(n,t) = cos( InternVideo2_text(action_n), InternVideo2_video(clip_t) )  — frozen, video-native, no classifier."),
    (0, "+ MASRA (Eliz):  LRCA / ESTA training-time losses that shape s toward the transcript's "
        "semantic-relational structure; MLLM at train only."),
    (0, "+ CVA:  CTE-style encoder on the VLM features; CBD boundary-contrastive loss on "
        "the aligned boundary frames."),
    (0, "Order-preserving alignment on s (ASOT — already in DELTA as --model_type wclot, ~1 line to "
        "swap the cost matrix)  →  Y*."),
    (0, "Y*  →  DELTA decoder + CRF + duration head (unchanged)  →  Obs%/Pred% MoC on 50Salads."),
    (0, "Contribution: VLM-direct transcript→frame alignment for weakly-supervised dense anticipation — unaddressed."),
], sub="InternVideo2 similarity  +  MASRA regulariser  +  CVA aligner  →  DELTA")

# ------------------------------------------------------------------ 18. SEQUENCING
table_slide("Sequencing — two VLM tracks, 50Salads",
    ["When", "MASRA track", "CVA track"],
    [
        ["Tomorrow", "M1: get MASRA code; reproduce on TACoS; read LRCA / ESTA / DAI", "get CVA code; reproduce on QVHighlights / TACoS"],
        ["+1–2 wk", "M2: internvideo2 backbone; chase raw 50S video; build s once available", "isolate CTE + CBD as reusable modules"],
        ["+2–3 wk", "M3–M4: adapt LRCA / ESTA to the transcript setting", "CTE / CBD on our VLM features (once s exists)"],
        ["+3–4 wk", "M5–M6: order-preserving alignment; score Y* (TA metrics) vs baselines", "join: CBD on the aligned boundaries"],
        ["later", "M7: Y* → DELTA decoder → Obs%/Pred% MoC", "ablations; write-up"],
    ],
    sub="both blocked on raw 50Salads video for the 50S runs; TACoS reproduction can start now",
    col_widths=[1.2, 5.7, 5.2], font=9.5)

# ------------------------------------------------------------------ 19. QUESTIONS
bullets("Questions for you", [
    (0, "Can the group share raw 50Salads video?  (blocks the VLM feature extraction for both of us.)"),
    (0, "MASRA uses an MLLM at training to generate captions — for us the transcript is already given. "
        "Keep the MLLM (to expand labels into descriptions), or drop it and use LRCA/ESTA with the plain labels?"),
    (0, "Primary target metric — alignment Y* quality on 50Salads, or downstream DLTA MoC?"),
    (0, "Is establishing a 50Salads transcript-only alignment benchmark a contribution?"),
    (0, "Frozen VLM only, or is fine-tuning the alignment cost (light adapter) in scope?"),
    (0, "VLM backbone — InternVideo2, or something the group already uses?"),
    (0, "Do we have the DELTA supplementary material (loss weights, decoder hyperparameters)?"),
], sub="to set the next phase")

# ------------------------------------------------------------------ CLOSE
s = section("Thank you", kicker="github.com/epayaslii/DELTA")
box = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.2))
tf = _tf(box)
r = tf.paragraphs[0].add_run()
r.text = ("Docs: temporal-alignment.md · 50salads-notes.md · delta-code.md · "
          "baselines-hal-cva.md · approach.md")
r.font.size = Pt(13); r.font.name = FONT; r.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE3)

out = "/Users/elizpayasli/Documents/GitHub/DELTA/slides/DELTA_progress.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
